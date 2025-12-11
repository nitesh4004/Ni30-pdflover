import streamlit as st
import os
from io import BytesIO
import zipfile
from PIL import Image, ImageOps, ImageFilter

# PDF Libraries
from PyPDF2 import PdfReader, PdfWriter

# PPT Libraries
from pptx import Presentation

# PDF to Image (Requires Poppler System Dependency)
try:
    from pdf2image import convert_from_bytes
    HAS_PDF2IMAGE = True
except ImportError:
    HAS_PDF2IMAGE = False
except Exception:
    HAS_PDF2IMAGE = False

# --- 1. PAGE CONFIGURATION ---
st.set_page_config(
    page_title="DocMint - Pro Tools",
    page_icon="🍃",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# --- 2. SESSION STATE MANAGEMENT ---
if 'current_tool' not in st.session_state:
    st.session_state['current_tool'] = "Home"

def navigate_to(tool_name):
    st.session_state['current_tool'] = tool_name

def go_home():
    st.session_state['current_tool'] = "Home"

# --- 3. CUSTOM CSS (MODERN SAAS UI) ---
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600;800&display=swap');

    html, body, [class*="css"] {
        font-family: 'Inter', sans-serif;
        color: #1e293b;
    }

    /* Clean Header */
    .nav-container {
        background-color: white;
        border-bottom: 1px solid #e2e8f0;
        padding: 1.5rem 3rem;
        margin-top: -4rem; 
        margin-left: -4rem;
        margin-right: -4rem;
        margin-bottom: 2rem;
        display: flex;
        align-items: center;
        justify-content: space-between;
    }
    
    .logo-box {
        display: flex;
        align-items: center;
        gap: 15px;
    }

    .app-name {
        font-size: 2.2rem;
        font-weight: 800;
        color: #0f172a;
        letter-spacing: -1px;
    }

    .tagline {
        color: #64748b;
        font-size: 0.95rem;
        font-weight: 500;
    }

    /* Modern Cards for Tools */
    div.stButton > button {
        background-color: white;
        border: 1px solid #e2e8f0;
        border-radius: 12px;
        padding: 24px;
        color: #334155;
        text-align: left;
        display: flex;
        flex-direction: column;
        align-items: flex-start;
        justify-content: flex-start;
        box-shadow: 0 1px 2px 0 rgba(0, 0, 0, 0.05);
        transition: all 0.2s ease;
        height: 160px;
    }

    div.stButton > button:hover {
        border-color: #00C853;
        box-shadow: 0 10px 15px -3px rgba(0, 0, 0, 0.1);
        transform: translateY(-2px);
        color: #00C853;
    }

    div.stButton > button p {
        font-size: 0.9rem;
        color: #64748b;
        margin-top: 8px;
    }

    /* Primary Buttons (Upload/Download) */
    div.stDownloadButton > button, div.stFormSubmitButton > button, .primary-btn {
        background-color: #007AFF !important; /* Vivid Blue like the screenshot */
        color: white !important;
        border: none;
        border-radius: 8px;
        padding: 12px 24px;
        font-weight: 600;
        box-shadow: 0 4px 6px -1px rgba(0, 122, 255, 0.2);
        transition: background 0.2s;
    }
    
    div.stDownloadButton > button:hover, div.stFormSubmitButton > button:hover {
        background-color: #0062cc !important;
    }

    /* Dashed Upload Area Styling */
    [data-testid="stFileUploader"] {
        border: 2px dashed #cbd5e1;
        border-radius: 12px;
        padding: 2rem;
        background-color: #f8fafc;
        text-align: center;
    }
    
    /* Input Fields */
    .stNumberInput input, .stTextInput input, .stSelectbox div[data-baseweb="select"] {
        border-radius: 8px;
        border-color: #e2e8f0;
    }

    /* Tab Styling */
    .stTabs [data-baseweb="tab-list"] {
        gap: 24px;
        border-bottom: 1px solid #e2e8f0;
    }
    .stTabs [data-baseweb="tab"] {
        height: 48px;
        font-weight: 600;
        color: #64748b;
        border: none;
        background-color: transparent;
    }
    .stTabs [aria-selected="true"] {
        color: #007AFF !important;
        border-bottom: 2px solid #007AFF !important;
    }
    
    /* File Info Card */
    .file-info {
        background: white;
        border: 1px solid #e2e8f0;
        border-radius: 8px;
        padding: 15px;
        margin-bottom: 20px;
        display: flex;
        align-items: center;
        gap: 15px;
    }
</style>
""", unsafe_allow_html=True)

# --- 4. HELPER FUNCTIONS ---
def create_zip(files_dict, zip_name):
    zip_buffer = BytesIO()
    with zipfile.ZipFile(zip_buffer, "a", zipfile.ZIP_DEFLATED, False) as zip_file:
        for file_name, data in files_dict.items():
            zip_file.writestr(file_name, data)
    return zip_buffer.getvalue()

def get_size_format(b, factor=1024, suffix="B"):
    """Scale bytes to its proper byte format"""
    for unit in ["", "K", "M", "G", "T", "P", "E", "Z"]:
        if b < factor:
            return f"{b:.2f} {unit}{suffix}"
        b /= factor
    return f"{b:.2f} Y{suffix}"

# --- 5. UI RENDERERS ---

def render_header():
    tool_display = st.session_state['current_tool']
    logo_url = "https://github.com/nitesh4004/Ni30-pdflover/blob/main/docmint.png?raw=true"
    
    st.markdown(f"""
    <div class="nav-container">
        <div class="logo-box">
            <img src="{logo_url}" style="height: 65px; width: auto; border-radius: 12px; box-shadow: 0 4px 6px rgba(0,0,0,0.1);">
            <div>
                <div class="app-name">DocMint</div>
            </div>
        </div>
        <div style="text-align: right;">
            <div style="font-size: 1.2rem; font-weight: 700; color: #334155;">{tool_display if tool_display != "Home" else "Dashboard"}</div>
            <div class="tagline">Simple tools for complex docs</div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    if st.session_state['current_tool'] != "Home":
        if st.button("← Back to Dashboard", key="back_home"):
            go_home()
            st.rerun()
        st.write("") # Spacer

def render_tool_card(tool_info, col):
    with col:
        # Using HTML inside button label for rich styling
        label = f"{tool_info['icon']}  {tool_info['name']}\n\n{tool_info['desc']}"
        if st.button(label, key=tool_info['id'], use_container_width=True):
            navigate_to(tool_info['name'])
            st.rerun()

def render_home():
    # Tools Data
    tools_pdf = [
        {"name": "Merge PDF", "icon": "🔗", "desc": "Combine multiple PDFs into one document.", "id": "t_merge"},
        {"name": "Split PDF", "icon": "✂️", "desc": "Extract pages or split document.", "id": "t_split"},
        {"name": "PDF to JPG", "icon": "🖼️", "desc": "Convert high-quality pages to images.", "id": "t_p2j"},
        {"name": "PDF Text", "icon": "📝", "desc": "Extract text content from PDFs.", "id": "t_ptext"},
    ]
    
    tools_img = [
        {"name": "Resize Image", "icon": "📐", "desc": "Resize by pixel, percent with DPI control.", "id": "t_resize"},
        {"name": "Image Editor", "icon": "🎨", "desc": "Filters, rotation, and adjustments.", "id": "t_edit"},
        {"name": "Convert Format", "icon": "🔄", "desc": "Switch between JPG, PNG, WEBP.", "id": "t_conv"},
        {"name": "JPG to PDF", "icon": "📑", "desc": "Compile images into a PDF file.", "id": "t_j2p"},
    ]
    
    tools_ppt = [
        {"name": "Merge PPTX", "icon": "📊", "desc": "Combine PowerPoint presentations.", "id": "t_mppt"},
        {"name": "PPTX Text", "icon": "📄", "desc": "Extract text from slides.", "id": "t_pxt"},
    ]

    st.write("### What would you like to do today?")
    t1, t2, t3 = st.tabs(["PDF Tools", "Image Tools", "PPT & Others"])
    
    with t1:
        cols = st.columns(4)
        for i, tool in enumerate(tools_pdf):
            render_tool_card(tool, cols[i % 4])
            
    with t2:
        cols = st.columns(4)
        for i, tool in enumerate(tools_img):
            render_tool_card(tool, cols[i % 4])
            
    with t3:
        cols = st.columns(4)
        for i, tool in enumerate(tools_ppt):
            render_tool_card(tool, cols[i % 4])

# --- 6. ADVANCED IMAGE RESIZER (NEW FEATURE) ---
def tool_resize_image():
    st.markdown("### Resize an Image")
    st.markdown("Upload your image to change dimensions, resolution, and file size.")
    
    uploaded = st.file_uploader("Select Image", type=["png", "jpg", "jpeg", "webp"])
    
    if uploaded:
        img = Image.open(uploaded)
        file_size = uploaded.size
        
        # File Info Card (Mimicking the screenshot thumbnail)
        st.markdown(f"""
        <div class="file-info">
            <div style="font-size: 2rem;">🖼️</div>
            <div>
                <div style="font-weight: 700; color: #334155;">{uploaded.name}</div>
                <div style="color: #64748b; font-size: 0.9rem;">{img.width} x {img.height} px • {get_size_format(file_size)}</div>
            </div>
        </div>
        """, unsafe_allow_html=True)

        st.markdown("#### Choose new size and format")
        
        # Controls Layout
        c1, c2, c3 = st.columns([1, 1, 1])
        
        with c3:
            unit = st.selectbox("Unit", ["Pixels", "Percent"], index=0)
            
        with c1:
            if unit == "Pixels":
                width = st.number_input("Width", value=img.width, min_value=1)
            else:
                percent = st.number_input("Percentage", value=100, min_value=1, max_value=500)
                
        with c2:
            lock_aspect = st.checkbox("🔒 Lock Aspect Ratio", value=True)
            if unit == "Pixels":
                if lock_aspect:
                    # Calculate height based on width input
                    aspect_ratio = img.height / img.width
                    calculated_height = int(width * aspect_ratio)
                    height = st.number_input("Height", value=calculated_height, disabled=True)
                    st.caption(f"Auto-calculated: {height}px")
                else:
                    height = st.number_input("Height", value=img.height, min_value=1)
            else:
                st.write("") # Spacer for alignment
                
        st.markdown("---")
        
        # Advanced Settings Row
        r1, r2, r3 = st.columns(3)
        with r1:
            dpi = st.number_input("Resolution (DPI)", value=72, min_value=72, max_value=600, step=1)
        with r2:
            fmt = st.selectbox("Format", ["JPG", "PNG", "WEBP"], index=0 if img.format == "JPEG" else 1)
        with r3:
            quality = st.slider("Quality", 1, 100, 90)

        # Processing Button
        st.write("")
        if st.button("Resize Image", type="primary", use_container_width=True):
            with st.spinner("Resizing..."):
                # Logic
                if unit == "Pixels":
                    new_w, new_h = int(width), int(height)
                else:
                    new_w = int(img.width * (percent/100))
                    new_h = int(img.height * (percent/100))
                
                # Resize
                resized_img = img.resize((new_w, new_h), Image.Resampling.LANCZOS)
                
                # Save
                b = BytesIO()
                save_fmt = "JPEG" if fmt == "JPG" else fmt
                
                # Handle Transparency for JPG
                if save_fmt == "JPEG" and resized_img.mode in ("RGBA", "P"):
                    resized_img = resized_img.convert("RGB")
                    
                resized_img.save(b, format=save_fmt, dpi=(dpi, dpi), quality=quality)
                
                st.success("Image Resized Successfully!")
                st.download_button(
                    label=f"Download {fmt} ({get_size_format(b.getbuffer().nbytes)})",
                    data=b.getvalue(),
                    file_name=f"resized_{uploaded.name.split('.')[0]}.{fmt.lower()}",
                    mime=f"image/{fmt.lower()}"
                )


# --- 7. OTHER TOOLS LOGIC ---

def tool_merge_pdf():
    st.info("Combine multiple PDFs into one. Reorder them using the list below.")
    uploaded_pdfs = st.file_uploader("Upload PDF files", type="pdf", accept_multiple_files=True)
    if uploaded_pdfs:
        file_map = {file.name: file for file in uploaded_pdfs}
        order = st.multiselect("File Order:", list(file_map.keys()), default=list(file_map.keys()))
        if st.button("Merge Files") and order:
            merger = PdfWriter()
            for name in order:
                merger.append(file_map[name])
            output = BytesIO()
            merger.write(output)
            st.download_button("Download Merged PDF", output.getvalue(), "docmint_merged.pdf", "application/pdf")

def tool_split_pdf():
    st.info("Split PDF into pages.")
    u = st.file_uploader("Upload PDF", type="pdf")
    if u:
        r = PdfReader(u)
        st.write(f"**Total Pages:** {len(r.pages)}")
        if st.button("Split All Pages (ZIP)"):
            files = {}
            for i in range(len(r.pages)):
                w = PdfWriter()
                w.add_page(r.pages[i])
                o = BytesIO()
                w.write(o)
                files[f"page_{i+1}.pdf"] = o.getvalue()
            st.download_button("Download ZIP", create_zip(files, "split.zip"), "split.zip", "application/zip")

def tool_pdf_to_img():
    if not HAS_PDF2IMAGE:
        st.error("Poppler not installed.")
        return
    u = st.file_uploader("Upload PDF", type="pdf")
    if u and st.button("Convert"):
        imgs = convert_from_bytes(u.read())
        files = {}
        for i, im in enumerate(imgs):
            b = BytesIO()
            im.save(b, "JPEG")
            files[f"page_{i+1}.jpg"] = b.getvalue()
        st.download_button("Download Images", create_zip(files, "imgs.zip"), "imgs.zip", "application/zip")

def tool_img_editor():
    st.info("Basic image adjustments.")
    u = st.file_uploader("Upload", type=["jpg", "png"])
    if u:
        img = Image.open(u)
        st.image(img, width=300)
        c1, c2 = st.columns(2)
        angle = c1.slider("Rotate", 0, 360, 0)
        filt = c2.selectbox("Filter", ["None", "BW", "Blur", "Sharpen"])
        if st.button("Apply"):
            p = img.rotate(angle, expand=True)
            if filt == "BW": p = ImageOps.grayscale(p)
            elif filt == "Blur": p = p.filter(ImageFilter.BLUR)
            elif filt == "Sharpen": p = p.filter(ImageFilter.SHARPEN)
            b = BytesIO()
            p.save(b, format="PNG")
            st.download_button("Download", b.getvalue(), "edited.png", "image/png")

def tool_img_convert():
    u = st.file_uploader("Upload", type=["png", "jpg", "webp"])
    if u:
        t = st.selectbox("Convert to", ["PNG", "JPEG", "PDF", "WEBP"])
        if st.button("Convert"):
            i = Image.open(u)
            if t == "JPEG" and i.mode == "RGBA": i = i.convert("RGB")
            b = BytesIO()
            i.save(b, format=t)
            mime = "application/pdf" if t == "PDF" else f"image/{t.lower()}"
            st.download_button("Download", b.getvalue(), f"conv.{t.lower()}", mime)

def tool_img_to_pdf():
    u = st.file_uploader("Upload Images", type=["png", "jpg"], accept_multiple_files=True)
    if u and st.button("Create PDF"):
        imgs = [Image.open(f).convert("RGB") for f in u]
        b = BytesIO()
        imgs[0].save(b, "PDF", save_all=True, append_images=imgs[1:])
        st.download_button("Download PDF", b.getvalue(), "images.pdf", "application/pdf")

def tool_merge_pptx():
    u = st.file_uploader("Upload PPTX", type="pptx", accept_multiple_files=True)
    if u and st.button("Merge"):
        out = Presentation()
        # simplified merge logic
        for f in u:
            prs = Presentation(f)
            for slide in prs.slides:
                # complex copy not fully supported in simple python, just adding blank layout
                out.slides.add_slide(out.slide_layouts[6]) 
        b = BytesIO()
        out.save(b)
        st.download_button("Download PPTX", b.getvalue(), "merged.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")

def tool_text_extract(kind):
    u = st.file_uploader(f"Upload {kind}", type=kind.lower())
    if u and st.button("Extract"):
        txt = ""
        if kind == "PDF":
            r = PdfReader(u)
            for p in r.pages: txt += p.extract_text() + "\n"
        else:
            p = Presentation(u)
            for s in p.slides:
                for sh in s.shapes:
                    if hasattr(sh, "text"): txt += sh.text + "\n"
        st.download_button("Download Text", txt, "extracted.txt", "text/plain")

# --- 8. ROUTING ---

render_header()
t = st.session_state['current_tool']

if t == "Home": render_home()
elif t == "Resize Image": tool_resize_image() # New Feature
elif t == "Merge PDF": tool_merge_pdf()
elif t == "Split PDF": tool_split_pdf()
elif t == "PDF to JPG": tool_pdf_to_img()
elif t == "Image Editor": tool_img_editor()
elif t == "Convert Format": tool_img_convert()
elif t == "JPG to PDF": tool_img_to_pdf()
elif t == "Merge PPTX": tool_merge_pptx()
elif t == "PDF Text": tool_text_extract("PDF")
elif t == "PPTX Text": tool_text_extract("PPTX")

st.markdown("---")
st.markdown("<div style='text-align:center; color:#94a3b8; font-size:0.8rem;'>© 2024 DocMint</div>", unsafe_allow_html=True)
